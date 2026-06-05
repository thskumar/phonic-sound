"""Extract slide narration text from PPTX shapes or OCR on embedded images."""

from __future__ import annotations

import io
import logging
import re
from pathlib import Path
from typing import List, Optional, TYPE_CHECKING

from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide

logger = logging.getLogger(__name__)

# Skip tiny icons/decorations from Google Slides exports
MIN_PICTURE_BYTES = 8_000

_ocr_engine = None
_ocr_unavailable = False


def _get_ocr_engine():
    """Lazy-load RapidOCR (pip install rapidocr-onnxruntime)."""
    global _ocr_engine, _ocr_unavailable
    if _ocr_unavailable:
        return None
    if _ocr_engine is not None:
        return _ocr_engine
    try:
        from rapidocr_onnxruntime import RapidOCR

        _ocr_engine = RapidOCR()
        logger.info("OCR engine ready (rapidocr-onnxruntime)")
        return _ocr_engine
    except ImportError:
        logger.warning(
            "OCR unavailable. Install: pip install rapidocr-onnxruntime"
        )
        _ocr_unavailable = True
        return None


def _clean_ocr_text(raw: str) -> str:
    """Normalize OCR output for English TTS."""
    if not raw:
        return ""
    text = raw.replace("\n", " ").replace("|", " ")
    # Drop fallback-renderer boilerplate if composite was OCR'd by mistake
    text = re.sub(r"Slide\s+\d+\s+of\s+\d+", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"No text on this slide", " ", text, flags=re.IGNORECASE)
    # English phonics: keep ASCII letters, digits, common punctuation
    text = re.sub(r"[^A-Za-z0-9\s.,!?'\"-]", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _ocr_image(image: Image.Image) -> str:
    engine = _get_ocr_engine()
    if engine is None:
        return ""
    try:
        result, _ = engine(image)
        if not result:
            return ""
        parts = [str(item[1]).strip() for item in result if item and len(item) > 1]
        return _clean_ocr_text(" ".join(parts))
    except Exception as exc:
        logger.warning("OCR failed: %s", exc)
        return ""


def _ocr_path(path: Path) -> str:
    if not path.is_file():
        return ""
    try:
        with Image.open(path) as im:
            return _ocr_image(im.convert("RGB"))
    except Exception as exc:
        logger.warning("Could not OCR %s: %s", path, exc)
        return ""


def _picture_shapes(slide: "Slide") -> List:
    shapes = []

    def collect(shape) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                collect(child)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes.append(shape)

    for shape in slide.shapes:
        collect(shape)
    return shapes


def ocr_text_from_slide(slide: "Slide", composite_path: Optional[Path] = None) -> str:
    """
    Read text burned into slide images (Google Slides / picture-only decks).
    Uses largest embedded images first, then the rendered composite PNG.
    """
    pictures = _picture_shapes(slide)
    pictures.sort(key=lambda s: len(s.image.blob), reverse=True)

    collected: List[str] = []
    for shape in pictures:
        if len(shape.image.blob) < MIN_PICTURE_BYTES:
            continue
        try:
            img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
        except Exception:
            continue
        text = _ocr_image(img)
        if text and text not in collected:
            collected.append(text)
        # Main slide art is usually the largest image; stop after good hit
        if text and len(shape.image.blob) > 50_000:
            break

    if collected:
        merged = _clean_ocr_text(" ".join(collected))
        if merged:
            return merged

    # Only use composite when slide has no large embedded images
    if composite_path and not pictures:
        return _ocr_path(composite_path)

    return ""


def extract_all_slide_texts(
    prs: "Presentation",
    pptx_texts: List[str],
    slide_image_paths: Optional[List[Path]] = None,
    use_ocr: bool = True,
) -> List[str]:
    """
    Build narration text per slide: PPTX text frames, then OCR if empty.
    """
    paths = slide_image_paths or []
    result: List[str] = []

    for i, slide in enumerate(prs.slides):
        text = pptx_texts[i].strip() if i < len(pptx_texts) else ""
        if not text and use_ocr:
            composite = paths[i] if i < len(paths) else None
            text = ocr_text_from_slide(slide, composite)
            if text:
                preview = text[:60] + ("..." if len(text) > 60 else "")
                logger.info("Slide %s OCR: %s", i + 1, preview)
            else:
                logger.warning("Slide %s: no text found (PPTX or OCR)", i + 1)
        result.append(text)

    return result
