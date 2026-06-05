"""Extract slides from PPTX as PNG images and plain text."""

from __future__ import annotations

import io
import logging
import platform
import shutil
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from text_extractor import extract_all_slide_texts
from utils import ensure_dir, slide_filename

logger = logging.getLogger(__name__)

# Default slide dimensions (16:9)
SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080
FALLBACK_BG = (255, 248, 240)
FALLBACK_TEXT = (45, 55, 72)
FALLBACK_ACCENT = (99, 102, 241)
MIN_VALID_PNG_BYTES = 8_000


@dataclass
class SlideData:
    """One slide's extracted assets."""

    index: int
    image_path: Path
    text: str = ""
    title: str = ""


@dataclass
class ProcessResult:
    """Output of PPTX processing."""

    slides: List[SlideData] = field(default_factory=list)
    method: str = "unknown"


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            parts.append(line)
    return "\n".join(parts)


def extract_text_from_pptx(pptx_path: Path) -> List[str]:
    """Read all slide text via python-pptx."""
    prs = Presentation(str(pptx_path))
    texts: List[str] = []
    for slide in prs.slides:
        chunks: List[str] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    t = _shape_text(child)
                    if t:
                        chunks.append(t)
            else:
                t = _shape_text(shape)
                if t:
                    chunks.append(t)
        texts.append("\n".join(chunks).strip())
    return texts


def _get_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                continue
    return ImageFont.load_default()


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> List[str]:
    words = text.replace("\n", " \n ").split()
    lines: List[str] = []
    current: List[str] = []
    for word in words:
        if word == "\n":
            if current:
                lines.append(" ".join(current))
                current = []
            lines.append("")
            continue
        trial = " ".join(current + [word]) if current else word
        bbox = draw.textbbox((0, 0), trial, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current.append(word)
        else:
            if current:
                lines.append(" ".join(current))
            current = [word]
    if current:
        lines.append(" ".join(current))
    return lines or [""]


def render_fallback_slide(
    text: str,
    output_path: Path,
    slide_number: int,
    total: int,
) -> Path:
    """Render a simple child-friendly slide when PowerPoint COM is unavailable."""
    img = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), FALLBACK_BG)
    draw = ImageDraw.Draw(img)

    # Header bar
    draw.rectangle([0, 0, SLIDE_WIDTH, 120], fill=FALLBACK_ACCENT)
    title_font = _get_font(48)
    body_font = _get_font(56)
    small_font = _get_font(32)

    header = f"Slide {slide_number} of {total}"
    draw.text((60, 36), header, fill=(255, 255, 255), font=title_font)

    display = text.strip() or "(No text on this slide)"
    margin = 100
    max_w = SLIDE_WIDTH - 2 * margin
    y = 200
    for line in _wrap_text(draw, display, body_font, max_w):
        if not line:
            y += 30
            continue
        draw.text((margin, y), line, fill=FALLBACK_TEXT, font=body_font)
        bbox = draw.textbbox((margin, y), line, font=body_font)
        y += (bbox[3] - bbox[1]) + 24

    # Embedded images from pptx shapes handled separately in full fallback
    img.save(output_path, "PNG", optimize=True)
    return output_path


def _extract_images_from_slide(slide, dest: Path, prefix: str) -> List[Path]:
    """Save picture shapes from a slide."""
    saved: List[Path] = []
    idx = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                ext = shape.image.ext or "png"
                out = dest / f"{prefix}_img{idx}.{ext}"
                out.write_bytes(blob)
                saved.append(out)
                idx += 1
            except Exception as exc:
                logger.warning("Could not extract image: %s", exc)
    return saved


def _image_has_content(img: Image.Image) -> bool:
    """True if the image is not a single flat (e.g. all-white) color."""
    extrema = img.convert("RGB").resize((64, 64)).getextrema()
    # extrema = ((rmin, rmax), (gmin, gmax), (bmin, bmax))
    return any(mn != mx for (mn, mx) in extrema)


def _is_valid_slide_image(path: Path) -> bool:
    """Reject tiny or blank exports from PowerPoint COM."""
    if not path.is_file() or path.stat().st_size < MIN_VALID_PNG_BYTES:
        return False
    try:
        with Image.open(path) as im:
            return _image_has_content(im)
    except Exception:
        return False


def _paste_shapes_on_canvas(canvas: Image.Image, shapes, sw: int, sh: int) -> None:
    """Paste picture shapes onto canvas using slide coordinates."""
    cw, ch = canvas.size

    def paste(shape) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                paste(child)
            return
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return
        try:
            pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
        except Exception as exc:
            logger.warning("Could not load picture shape: %s", exc)
            return
        x = int(shape.left * cw / sw)
        y = int(shape.top * ch / sh)
        w = max(1, int(shape.width * cw / sw))
        h = max(1, int(shape.height * ch / sh))
        pic = pic.resize((w, h), Image.Resampling.LANCZOS)
        canvas.paste(pic, (x, y), pic)

    for shape in shapes:
        paste(shape)


def render_slide_composite(slide, prs: Presentation, output_path: Path) -> Path:
    """Build slide PNG by compositing embedded images at their layout positions."""
    sw, sh = prs.slide_width, prs.slide_height
    canvas = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), (255, 255, 255))
    _paste_shapes_on_canvas(canvas, slide.shapes, sw, sh)
    canvas.save(output_path, "PNG", optimize=True)
    return output_path


def _collect_pictures(shapes) -> List:
    """Recursively gather picture shapes (including inside groups)."""
    pics: List = []

    def walk(shape) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                walk(child)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics.append(shape)

    for shape in shapes:
        walk(shape)
    return pics


def _shape_area_fraction(shape, sw: int, sh: int) -> float:
    try:
        return ((shape.width or 0) * (shape.height or 0)) / float(sw * sh)
    except Exception:
        return 0.0


def render_clean_text_slide(text: str, output_path: Path) -> Path:
    """Child-friendly slide with centered text and no header/footer captions."""
    img = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), FALLBACK_BG)
    draw = ImageDraw.Draw(img)
    display = (text or "").strip()
    if display:
        font = _get_font(72)
        margin = 140
        max_w = SLIDE_WIDTH - 2 * margin
        lines = _wrap_text(draw, display, font, max_w)
        line_h = 96
        visible = [ln for ln in lines if ln != ""]
        total_h = line_h * max(1, len(visible))
        y = max(margin, (SLIDE_HEIGHT - total_h) // 2)
        for line in lines:
            if not line:
                y += line_h // 2
                continue
            bbox = draw.textbbox((0, 0), line, font=font)
            w = bbox[2] - bbox[0]
            draw.text(((SLIDE_WIDTH - w) // 2, y), line, fill=FALLBACK_TEXT, font=font)
            y += line_h
    img.save(output_path, "PNG", optimize=True)
    return output_path


def render_slide_image(
    slide,
    prs: Presentation,
    output_path: Path,
    text: str = "",
    background_threshold: float = 0.85,
) -> Path:
    """
    Render a clean slide PNG.

    Google Slides decks store each page as one full-bleed background image
    plus small decorative overlays (navigation arrows, sparkles, speaker
    icons). We render only the dominant background image, so those overlays
    are dropped and no 'Slide X of Y' / 'no text' captions are ever added.
    Falls back to a faithful composite, then to a clean text slide.
    """
    sw, sh = prs.slide_width, prs.slide_height
    pics = _collect_pictures(slide.shapes)

    if pics:
        biggest = max(pics, key=lambda s: _shape_area_fraction(s, sw, sh))
        if _shape_area_fraction(biggest, sw, sh) >= background_threshold:
            try:
                bg = Image.open(io.BytesIO(biggest.image.blob)).convert("RGB")
                bg = bg.resize((SLIDE_WIDTH, SLIDE_HEIGHT), Image.Resampling.LANCZOS)
                bg.save(output_path, "PNG", optimize=True)
                return output_path
            except Exception as exc:
                logger.warning("Background render failed: %s", exc)

        # No full-bleed background: faithfully composite all shapes
        canvas = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), FALLBACK_BG)
        _paste_shapes_on_canvas(canvas, slide.shapes, sw, sh)
        if _image_has_content(canvas):
            canvas.save(output_path, "PNG", optimize=True)
            return output_path

    return render_clean_text_slide(text, output_path)


def _composite_images_on_slide(base_path: Path, image_paths: List[Path]) -> None:
    """Overlay extracted images on fallback slide (simple grid)."""
    if not image_paths:
        return
    base = Image.open(base_path).convert("RGBA")
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    x, y = 100, SLIDE_HEIGHT // 2
    for img_path in image_paths[:3]:
        try:
            pic = Image.open(img_path).convert("RGBA")
            pic.thumbnail((600, 400), Image.Resampling.LANCZOS)
            overlay.paste(pic, (x, y), pic)
            x += pic.width + 40
        except Exception as exc:
            logger.warning("Overlay failed for %s: %s", img_path, exc)
    combined = Image.alpha_composite(base, overlay)
    combined.convert("RGB").save(base_path, "PNG", optimize=True)


def _export_via_win32com(pptx_path: Path, slides_dir: Path) -> Optional[List[Path]]:
    """Export slides using pywin32 (most reliable on Windows)."""
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return None

    pptx_path = pptx_path.resolve()
    slides_dir = ensure_dir(slides_dir)
    paths: List[Path] = []
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(
            str(pptx_path), WithWindow=False, ReadOnly=True
        )
        count = presentation.Slides.Count
        for i in range(1, count + 1):
            out = slides_dir / slide_filename(i - 1, "png")
            base = str((slides_dir / f"_ppt_export_{i}").resolve())
            presentation.Slides(i).Export(base, "PNG", SLIDE_WIDTH, SLIDE_HEIGHT)
            time.sleep(0.05)
            exported = _find_exported_png(slides_dir, f"_ppt_export_{i}")
            if exported and exported != out:
                shutil.move(str(exported), str(out))
            elif not out.is_file() and exported:
                shutil.move(str(exported), str(out))
            paths.append(out)
            logger.debug("Exported slide %s via win32com", i)
        return paths if all(_is_valid_slide_image(p) for p in paths) else None
    except Exception as exc:
        logger.warning("win32com PowerPoint export failed: %s", exc)
        return None
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass


def _export_via_comtypes(pptx_path: Path, slides_dir: Path) -> Optional[List[Path]]:
    """Export slides using comtypes EnsureDispatch."""
    try:
        import comtypes.client  # type: ignore
    except ImportError:
        return None

    pptx_path = pptx_path.resolve()
    slides_dir = ensure_dir(slides_dir)
    paths: List[Path] = []
    powerpoint = None
    presentation = None
    try:
        powerpoint = comtypes.client.gencache.EnsureDispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        count = presentation.Slides.Count
        for i in range(1, count + 1):
            out = slides_dir / slide_filename(i - 1, "png")
            base = str((slides_dir / f"_ppt_export_{i}").resolve())
            presentation.Slides(i).Export(base, "PNG", SLIDE_WIDTH, SLIDE_HEIGHT)
            time.sleep(0.05)
            exported = _find_exported_png(slides_dir, f"_ppt_export_{i}")
            if exported and exported != out:
                shutil.move(str(exported), str(out))
            paths.append(out)
            logger.debug("Exported slide %s via comtypes", i)
        return paths if all(_is_valid_slide_image(p) for p in paths) else None
    except Exception as exc:
        logger.warning("comtypes PowerPoint export failed: %s", exc)
        return None
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass


def _find_exported_png(slides_dir: Path, base_name: str) -> Optional[Path]:
    """Locate PNG created by PowerPoint Export (may vary in extension)."""
    for ext in (".png", ".PNG"):
        candidate = slides_dir / f"{base_name}{ext}"
        if candidate.is_file():
            return candidate
    matches = list(slides_dir.glob(f"{base_name}*.png")) + list(
        slides_dir.glob(f"{base_name}*.PNG")
    )
    return matches[0] if matches else None


def export_via_com(pptx_path: Path, slides_dir: Path) -> Optional[List[Path]]:
    """
    Export slides as PNG using Microsoft PowerPoint COM (Windows only).
    Returns list of image paths or None if unavailable.
    """
    if platform.system() != "Windows":
        logger.info("COM export skipped: not on Windows")
        return None

    paths = _export_via_win32com(pptx_path, slides_dir)
    if paths:
        return paths
    return _export_via_comtypes(pptx_path, slides_dir)


def process_pptx(
    pptx_path: Path,
    slides_dir: Path,
    force_fallback: bool = False,
    use_ocr: bool = True,
) -> ProcessResult:
    """
    Extract slides to PNG and collect text.
    Prefers COM on Windows; falls back to python-pptx + Pillow rendering.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.is_file():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    slides_dir = ensure_dir(slides_dir)
    texts = extract_text_from_pptx(pptx_path)
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    if total == 0:
        raise ValueError("Presentation has no slides")

    image_paths: Optional[List[Path]] = None
    method = "fallback"

    slide_data_list: List[SlideData] = []
    image_paths: Optional[List[Path]] = None
    method = "composite"

    if not force_fallback:
        image_paths = export_via_com(pptx_path, slides_dir)
        if image_paths and len(image_paths) == total:
            method = "powerpoint_com"
            logger.info("Exported %s slides via PowerPoint COM", total)
        else:
            image_paths = None
            if not force_fallback:
                logger.info(
                    "PowerPoint COM export missing or blank; using image composite renderer"
                )

    if image_paths and len(image_paths) == total:
        for i, img_path in enumerate(image_paths):
            slide_data_list.append(
                SlideData(
                    index=i,
                    image_path=img_path,
                    text=texts[i] if i < len(texts) else "",
                    title=_first_line(texts[i] if i < len(texts) else ""),
                )
            )
    else:
        # Clean renderer: full-bleed background image, decorations/captions dropped
        method = "image"
        for i, slide in enumerate(prs.slides):
            out = slides_dir / slide_filename(i, "png")
            text = texts[i] if i < len(texts) else ""
            render_slide_image(slide, prs, out, text)
            slide_data_list.append(
                SlideData(
                    index=i,
                    image_path=out,
                    text=text,
                    title=_first_line(text),
                )
            )
        logger.info("Rendered %s slides (clean image renderer)", total)

    # Google Slides decks store text inside images — fill narration via OCR
    image_paths_list = [sd.image_path for sd in slide_data_list]
    enriched = extract_all_slide_texts(prs, texts, image_paths_list, use_ocr=use_ocr)
    for i, sd in enumerate(slide_data_list):
        sd.text = enriched[i]
        sd.title = _first_line(enriched[i])

    ocr_count = sum(1 for t in enriched if t.strip())
    logger.info("Narration text ready for %s/%s slides", ocr_count, total)
    return ProcessResult(slides=slide_data_list, method=method)


def _first_line(text: str) -> str:
    for line in text.splitlines():
        line = line.strip()
        if line:
            return line
    return ""
