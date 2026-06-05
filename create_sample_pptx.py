"""Create a minimal sample lesson.pptx for testing the pipeline."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def main() -> None:
    prs = Presentation()
    slides_content = [
        ("Welcome!", "Hello! Today we will learn something fun."),
        ("Counting", "One, two, three. Can you count with me?"),
        ("Well Done!", "Great job! You finished the lesson."),
    ]
    for title, body in slides_content:
        layout = prs.slide_layouts[1]  # title and content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        for para in slide.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(28)

    out = Path("lesson.pptx")
    prs.save(str(out))
    print(f"Created {out.resolve()}")


if __name__ == "__main__":
    main()
